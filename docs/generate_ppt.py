import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme (Modern Premium Deep Blue & Steel Gray)
    c_primary = RGBColor(0x1F, 0x4E, 0x79)   # Deep Blue
    c_secondary = RGBColor(0x2F, 0x55, 0x97) # Steel Blue
    c_text_dark = RGBColor(0x2B, 0x2B, 0x2B) # Soft Black
    c_text_light = RGBColor(0x59, 0x59, 0x59)# Gray
    c_accent = RGBColor(0xD9, 0x53, 0x4F)    # Crimson/Amber warning accent
    
    # Helper to set slide background
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to format a title text block
    def format_title(title_p, text, font_size=36, color=c_primary):
        title_p.text = text
        title_p.font.name = 'Calibri'
        title_p.font.size = Pt(font_size)
        title_p.font.bold = True
        title_p.font.color.rgb = color
        title_p.alignment = PP_ALIGN.LEFT

    # Helper to add standard bullets
    def add_bullet_point(tf, bold_pref, text_body, color=c_text_dark, size=16, level=0):
        p = tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)
        
        # Bold Prefix
        run_bold = p.add_run()
        run_bold.text = bold_pref
        run_bold.font.name = 'Calibri'
        run_bold.font.bold = True
        run_bold.font.size = Pt(size)
        run_bold.font.color.rgb = c_secondary
        
        # Text Body
        run_text = p.add_run()
        run_text.text = text_body
        run_text.font.name = 'Calibri'
        run_text.font.bold = False
        run_text.font.size = Pt(size)
        run_text.font.color.rgb = color

    # -----------------------------------------------------------------------
    # SLIDE 1: Cover Page (Custom Blank Layout for maximum control)
    # -----------------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank Slide
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Accent shape on the left border (Visual polish)
    left_bar = slide1.shapes.add_shape(
        1, # Rectangle
        Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = c_primary
    left_bar.line.fill.background() # No border
    
    # Title Textbox
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "MONEY MANAGER JAVA SUITE"
    p1.font.name = 'Calibri Light'
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = c_primary
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "JavaFX Desktop & Spring Boot Web Applications"
    p1_sub.font.name = 'Calibri'
    p1_sub.font.size = Pt(20)
    p1_sub.font.italic = True
    p1_sub.font.color.rgb = c_secondary
    p1_sub.space_before = Pt(8)
    
    # Subtitle Details
    details_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.0), Inches(1.5))
    tf_det = details_box.text_frame
    p_det1 = tf_det.paragraphs[0]
    p_det1.text = "System Architecture & SOLID Requirements Specification"
    p_det1.font.name = 'Calibri'
    p_det1.font.size = Pt(14)
    p_det1.font.bold = True
    p_det1.font.color.rgb = c_text_dark
    
    p_det2 = tf_det.add_paragraph()
    p_det2.text = "Platform: Java 21 LTS  •  Databases: SQLite / PostgreSQL\nDocument Version: 1.0  •  Status: Approved"
    p_det2.font.name = 'Calibri'
    p_det2.font.size = Pt(12)
    p_det2.font.color.rgb = c_text_light
    p_det2.space_before = Pt(4)

    # -----------------------------------------------------------------------
    # SLIDE 2: Project Overview
    # -----------------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    format_title(slide2.shapes.title.text_frame.paragraphs[0], "1. Project Overview & Scope")
    
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    add_bullet_point(tf2, "Unified Personal Finance Suite: ", "Synchronizes desktop workstation utility with multi-screen enterprise web capabilities.")
    add_bullet_point(tf2, "JavaFX Desktop Client: ", "Standalone, offline-first application running on Java 21 LTS with data persisted in embedded SQLite database engines.")
    add_bullet_point(tf2, "Spring Boot Web Client: ", "High-concurrency web application powered by Java 21, Spring Boot 4.0.6, Thymeleaf layouts, and PostgreSQL.")
    add_bullet_point(tf2, "Core Capabilities: ", "Secure user profiles, transaction logs, relational category budgets with amber/red progress warnings, savings milestones, and notes lists.")

    # -----------------------------------------------------------------------
    # SLIDE 3: 3-Tier System Architecture
    # -----------------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide3.shapes.title.text_frame.paragraphs[0], "2. 3-Tier Architecture Boundaries")
    
    content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf3 = content_box3.text_frame
    tf3.word_wrap = True
    
    add_bullet_point(tf3, "Presentation Tier: ", "JavaFX FXML layouts & controllers (Desktop) and Spring WebMVC controllers mapping reactive Thymeleaf views (Web).")
    add_bullet_point(tf3, "Business Service Tier: ", "Core services (TransactionService, BudgetService, AuthService) managing rules, password checks, and validations, independent of GUI structures.")
    add_bullet_point(tf3, "Relational Repository Tier: ", "JDBC-based database repositories (JdbcTransactionRepo, JdbcUserRepo) executing statements directly to target databases.")
    add_bullet_point(tf3, "Interchangeable Connections: ", "Repositories map to SQLite local drivers (`money-manager.db`) or multi-user PostgreSQL services natively via standard interface layers.")

    # -----------------------------------------------------------------------
    # SLIDE 4: SOLID Design Principles
    # -----------------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide4.shapes.title.text_frame.paragraphs[0], "3. SOLID Design Implementation")
    
    content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True
    
    add_bullet_point(tf4, "Single Responsibility (SRP): ", "Classes have one focus (e.g. AuthService parses credentials; JdbcBudgetRepo queries budget DB rows).")
    add_bullet_point(tf4, "Open / Closed (OCP): ", "New reporting utilities or database connectors are added by declaring new classes, without editing core engines.")
    add_bullet_point(tf4, "Liskov Substitution (LSP): ", "JdbcTransactionRepo implements ITransactionRepo, meaning it can be replaced by MockTransactionRepo without breaking services.")
    add_bullet_point(tf4, "Interface Segregation (ISP): ", "Granular interfaces split data bounds: ITransactionRepo, IBudgetRepo, IGoalRepo, INoteRepo.")
    add_bullet_point(tf4, "Dependency Inversion (DIP): ", "Services depend on abstractions (interfaces), wired at runtime via constructor dependency injection (enables JUnit test coverage).")

    # -----------------------------------------------------------------------
    # SLIDE 5: Unified Relational Schema
    # -----------------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide5.shapes.title.text_frame.paragraphs[0], "4. Unified Relational Schema Design")
    
    content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf5 = content_box5.text_frame
    tf5.word_wrap = True
    
    add_bullet_point(tf5, "Table: users — ", "Saves secure profile coordinates. Passwords are encrypted using modern BCrypt hashing routines.")
    add_bullet_point(tf5, "Table: transactions — ", "Registers income and expense entries with title, positive amount (> 0), category tag, and date.")
    add_bullet_point(tf5, "Table: budgets — ", "Configures category-specific monthly spending caps. Enforces unique composite limits on user, category, and date.")
    add_bullet_point(tf5, "Table: savings_goals & goal_contributions — ", "Tracks milestone requirements, deadline alerts, contributions, and percentage completion bars.")
    add_bullet_point(tf5, "Table: user_settings & monthly_balance — ", "Tracks monthly target income indices and saves calculated monthly cash flows per profile.")

    # -----------------------------------------------------------------------
    # SLIDE 6: Key Features & Validation
    # -----------------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide6.shapes.title.text_frame.paragraphs[0], "5. Features & Validation Logic")
    
    content_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf6 = content_box6.text_frame
    tf6.word_wrap = True
    
    add_bullet_point(tf6, "Secure User Authentication: ", "Desktop security answers and web password verifications execute safely using BCrypt wrappers.")
    add_bullet_point(tf6, "Real-Time Spending Indicators: ", "UI progress bars automatically highlight alerts: green (<80%), amber warning (80-99%), and red indicator (>=100% budget breach).")
    add_bullet_point(tf6, "Transaction Input Validations: ", "Service layer ensures transaction amounts are positive and dates do not occur in the future.")
    add_bullet_point(tf6, "Visual Graphical Analytics: ", "Pie charts render categorical expense breakdowns and Bar charts display monthly net cash flow trends.")

    # -----------------------------------------------------------------------
    # SLIDE 7: Future Development Roadmap
    # -----------------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide7.shapes.title.text_frame.paragraphs[0], "6. Future Development Roadmap")
    
    content_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf7 = content_box7.text_frame
    tf7.word_wrap = True
    
    add_bullet_point(tf7, "Cross-Platform Mobile Clients: ", "Build reactive mobile frontends using Flutter or React Native mapping to shared backend business API controllers.")
    add_bullet_point(tf7, "Cloud Backup Synchronizer: ", "Introduce connection layers securely pushing local SQLite records to hosted AWS RDS PostgreSQL databases.")
    add_bullet_point(tf7, "AI-Powered Budget Coaching: ", "Deploy localized ML models evaluating transactions to suggest savings goals and flag anomalous expenditures.")
    add_bullet_point(tf7, "Document Export Engine: ", "Add PDF bank statements generated via OpenPDF and pre-formulated Excel reports compiled via Apache POI.")

    # Save presentation
    output_path = "docs/Money_Manager_Presentation.pptx"
    prs.save(output_path)
    print(f"[SUCCESS] Presentation generated and saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
